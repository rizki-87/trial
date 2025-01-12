from concurrent.futures import ThreadPoolExecutor      
import streamlit as st      
import tempfile      
from pathlib import Path      
from pptx import Presentation      
from spellchecker import SpellChecker      
import language_tool_python      
import csv      
import re      
import string      
from pptx.dml.color import RGBColor      
import logging      
import time
import pandas as pd  
from pydantic import BaseModel      
from utils.highlight import highlight_ppt, save_to_csv      
from utils.font_validation import validate_fonts_slide      
from utils.grammar_validation import initialize_language_tool, validate_grammar_slide      
from utils.spelling_validation import is_exempted, validate_spelling_slide      
from utils.decimal_validation import validate_decimal_consistency      
from utils.million_notation_validation import validate_million_notations_with_pandas      
from utils.validation import validate_tables, validate_charts      
from config import PREDEFINED_PASSWORD, TECHNICAL_TERMS, NUMERIC_TERMS      
  
# Initialize LanguageTool      
grammar_tool = initialize_language_tool()      
  
# Initialize SpellChecker      
spell = SpellChecker()      
spell.word_frequency.load_words(TECHNICAL_TERMS.union(NUMERIC_TERMS))      
  
# Configure logging      
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')      
  
# Password Protection      
def password_protection():      
    if "authenticated" not in st.session_state:      
        st.session_state.authenticated = False      
    if not st.session_state.authenticated:      
        with st.form("password_form", clear_on_submit=True):      
            password_input = st.text_input("Enter Password", type="password")      
            submitted = st.form_submit_button("Submit")      
            if submitted and password_input == PREDEFINED_PASSWORD:      
                st.session_state.authenticated = True      
                st.success("Access Granted! Please click 'Submit' again to proceed.")      
            elif submitted:      
                st.error("Incorrect Password")      
        return False      
    return True      
  
def extract_text_from_ppt(ppt_file):  
    text_data = []  
    presentation = Presentation(ppt_file)  
  
    for slide in presentation.slides:  
        slide_text = []  
        for shape in slide.shapes:  
            if shape.has_text_frame:  
                for paragraph in shape.text_frame.paragraphs:  
                    for run in paragraph.runs:  
                        slide_text.append(run.text)  
        text_data.append(" ".join(slide_text))  # Gabungkan teks dalam satu slide  
  
    return pd.DataFrame(text_data, columns=['text'])  # Kembalikan DataFrame  
  
def validate_slide(slide, slide_index, default_font, spell, grammar_tool, decimal_places, selected_notation):      
    slide_issues = []      
    start_time = time.time()      
  
    # Validate Spelling      
    slide_issues.extend(validate_spelling_slide(slide, slide_index + 1))      
    # Validate Fonts      
    slide_issues.extend(validate_fonts_slide(slide, slide_index + 1, default_font))      
    # Validate Grammar      
    slide_issues.extend(validate_grammar_slide(slide, slide_index + 1, grammar_tool))      
    # Validate Decimal Consistency      
    slide_issues.extend(validate_decimal_consistency(slide, slide_index + 1, decimal_places))      
  
    # Mengambil teks dari slide untuk validasi notasi juta  
    slide_text = []  
    for shape in slide.shapes:  
        if shape.has_text_frame:  
            for paragraph in shape.text_frame.paragraphs:  
                for run in paragraph.runs:  
                    slide_text.append(run.text)  
    df = pd.DataFrame(slide_text, columns=['text'])  # Membuat DataFrame dari teks slide  
  
    # Validasi Notasi Juta  
    slide_issues.extend(validate_million_notations_with_pandas(df, selected_notation))  # Memasukkan selected_notation      
  
    elapsed_time = time.time() - start_time      
    logging.debug(f"Slide {slide_index + 1} validation completed in {elapsed_time:.2f} seconds.")      
  
    return slide_issues    
  
  
def main():      
    if not password_protection():      
        return      
  
    st.title("PPT Validator")      
    uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"])      
    font_options = ["Arial", "Calibri", "Times New Roman", "Verdana", "Helvetica", "EYInterstate"]      
    default_font = st.selectbox("Select the default font for validation", font_options)      
    decimal_places = st.number_input("Enter the number of decimal places for validation", min_value=0, max_value=10, value=1)      
    # Pilihan Notasi      
    notation_options = ["m", "M", "Mn"]      
    selected_notation = st.selectbox("Select Notation for Validation", notation_options)      
  
    validation_option = st.radio("Validation Option:", ["All Slides", "Custom Range"])      
  
    if uploaded_file:      
        with tempfile.TemporaryDirectory() as tmpdir:      
            temp_ppt_path = Path(tmpdir) / "uploaded_ppt.pptx"      
            with open(temp_ppt_path, "wb") as f:      
                f.write(uploaded_file.getbuffer())      
  
            presentation = Presentation(temp_ppt_path)      
            total_slides = len(presentation.slides)      
  
            # Rentang Slide      
            start_slide, end_slide = 1, total_slides      
            if validation_option == "Custom Range":      
                start_slide = st.number_input("From Slide", min_value=1, max_value=total_slides, value=1)      
                end_slide_default = min(total_slides, 100)      
                end_slide = st.number_input("To Slide", min_value=start_slide, max_value=total_slides, value=end_slide_default)      
  
            if st.button("Run Validation"):      
                progress_bar = st.progress(0)      
                progress_text = st.empty()      
                issues = []      
  
                # Proses Paralel      
                with ThreadPoolExecutor() as executor:      
                    futures = []      
                    for slide_index in range(start_slide - 1, end_slide):      
                        slide = presentation.slides[slide_index]      
                        futures.append(executor.submit(validate_slide, slide, slide_index, default_font, spell, grammar_tool, decimal_places, selected_notation))      
  
                    for i, future in enumerate(futures):      
                        slide_issues = future.result()      
                        issues.extend(slide_issues)      
                        progress_percent = int((i + 1) / len(futures) * 100)      
                        progress_text.text(f"Progress: {progress_percent}%")      
                        progress_bar.progress(progress_percent / 100)      
  
                # Simpan Hasil      
                csv_output_path = Path(tmpdir) / "validation_report.csv"      
                highlighted_ppt_path = Path(tmpdir) / "highlighted_presentation.pptx"      
                save_to_csv(issues, csv_output_path)      
                highlight_ppt(temp_ppt_path, highlighted_ppt_path, issues)      
  
                # Simpan hasil di session state      
                st.session_state['csv_output'] = csv_output_path.read_bytes()      
                st.session_state['ppt_output'] = highlighted_ppt_path.read_bytes()      
                st.session_state['validation_completed'] = True      
                st.session_state['issues'] = issues      
                st.session_state['log_output_path'] = str(Path(tmpdir) / "validation_log.txt")      
                st.success("Validation completed!")      
  
                # Tulis Log      
                log_output_path = st.session_state['log_output_path']      
                with open(log_output_path, "w") as log_file:      
                    for handler in logging.root.handlers[:]:      
                        logging.root.removeHandler(handler)      
                    logging.basicConfig(filename=log_output_path, level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')      
                    logging.debug(f"Validation completed with {len(issues)} issues.")      
                    for issue in issues:      
                        logging.debug(f"Issue: {issue}")      
  
                # Tampilkan Log di Streamlit      
                with open(log_output_path, "r") as log_file:      
                    log_content = log_file.read()      
                    st.text_area("Validation Log", value=log_content, height=300)  # Menampilkan log di antarmuka pengguna    
  
    # Tampilkan Tombol Unduh jika validasi telah selesai      
    if st.session_state.get('validation_completed', False):      
        if 'csv_output' in st.session_state:      
            st.download_button("Download Validation Report (CSV)", st.session_state['csv_output'], file_name="validation_report.csv")      
        if 'ppt_output' in st.session_state:      
            st.download_button("Download Highlighted PPT", st.session_state['ppt_output'], file_name="highlighted_presentation.pptx")      
  
if __name__ == "__main__":      
    main()   




