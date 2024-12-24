# from concurrent.futures import ThreadPoolExecutor
# import streamlit as st
# import tempfile
# from pathlib import Path
# from pptx import Presentation
# from spellchecker import SpellChecker
# import language_tool_python
# import csv
# import re
# import string
# from pptx.dml.color import RGBColor
# import logging

# # Initialize LanguageTool
# def initialize_language_tool():
#     try:
#         return language_tool_python.LanguageToolPublicAPI('en-US')
#     except Exception as e:
#         st.error(f"LanguageTool initialization failed: {e}")
#         return None

# grammar_tool = initialize_language_tool()

# # Custom dictionary
# TECHNICAL_TERMS = {
#     "TensorFlow", "Caret", "ML", "DALL-E", "MLOps", "PyTorch", "MENA", "Keras", "Scikit-learn", "NumPy", "Pandas",
#     "Matplotlib", "OpenAI", "GPT-3", "Deep Learning", "Neural Network", "Data Science", "Seaborn", "Jupyter",
#     "Anaconda", "Reinforcement Learning", "Supervised Learning", "Unsupervised Learning", "Natural Language Processing",
#     "Computer Vision", "Big Data", "Data Mining", "Feature Engineering", "Hyperparameter", "Gradient Descent",
#     "Convolutional Neural Network", "Recurrent Neural Network", "Support Vector Machine", "Decision Tree",
#     "Random Forest", "Ensemble Learning", "Clustering", "Dimensionality Reduction", "Principal Component Analysis",
#     "Exploratory Data Analysis", "Model Evaluation", "Cross-Validation", "Overfitting", "Underfitting",
#     "Batch Normalization", "Dropout", "Activation Function", "Loss Function", "Backpropagation", "Transfer Learning",
#     "Generative Adversarial Network", "Autoencoder", "Tokenization", "Embedding", "Word2Vec", "BERT", "OpenCV",
#     "Flask", "Django", "REST API", "GraphQL", "SQL", "NoSQL", "MongoDB", "PostgreSQL", "MySQL", "Firebase",
#     "Cloud Computing", "AWS", "Azure", "Google Cloud", "Docker", "Kubernetes", "CI/CD", "DevOps", "Agile", "Scrum",
#     "Kanban", "Git", "GitHub", "Bitbucket", "Version Control", "API", "SDK", "Microservices", "Blockchain",
#     "Cryptocurrency", "IoT", "Edge Computing", "Quantum Computing", "Augmented Reality", "Virtual Reality",
#     "3D Printing", "Cybersecurity", "Penetration Testing", "Phishing", "Malware", "Ransomware", "Firewall", "VPN",
#     "SSL", "Encryption", "Decryption", "Hashing", "Digital Signature", "Data Privacy", "GDPR"
# }

# NUMERIC_TERMS = {f"{i}+" for i in range(1, 101)}

# # Initialize SpellChecker
# spell = SpellChecker()
# spell.word_frequency.load_words(TECHNICAL_TERMS.union(NUMERIC_TERMS))

# # Exemption function
# def is_exempted(word):
#     return word in TECHNICAL_TERMS or re.match(r"^\d+\+?$", word)

# # Spelling Validation
# def validate_spelling_slide(slide, slide_index):
#     issues = []
#     for shape in slide.shapes:
#         if shape.has_text_frame:
#             for paragraph in shape.text_frame.paragraphs:
#                 for run in paragraph.runs:
#                     words = re.findall(r"\b[\w+]+\b", run.text)
#                     for word in words:
#                         clean_word = word.strip(string.punctuation)
#                         if is_exempted(clean_word):
#                             continue
#                         if clean_word.lower() not in spell:
#                             correction = spell.correction(clean_word)
#                             if correction and correction != clean_word:
#                                 issues.append({
#                                     'slide': slide_index,
#                                     'issue': 'Misspelling',
#                                     'text': word,
#                                     'corrected': correction
#                                 })
#     return issues

# # Font Validation
# def validate_fonts_slide(slide, slide_index, default_font):
#     issues = []
#     for shape in slide.shapes:
#         if shape.has_text_frame:
#             for paragraph in shape.text_frame.paragraphs:
#                 for run in paragraph.runs:
#                     if run.text.strip() and run.font.name != default_font:
#                         issues.append({
#                             'slide': slide_index,
#                             'issue': 'Inconsistent Font',
#                             'text': run.text,
#                             'corrected': f"Expected: {default_font}, Found: {run.font.name}"
#                         })
#     return issues

# # Grammar Validation
# def validate_grammar_slide(slide, slide_index):
#     issues = []
#     for shape in slide.shapes:
#         if shape.has_text_frame:
#             for paragraph in shape.text_frame.paragraphs:
#                 for run in paragraph.runs:
#                     text = run.text.strip()
#                     if text and grammar_tool:
#                         matches = grammar_tool.check(text)
#                         for match in matches:
#                             issues.append({
#                                 'slide': slide_index,
#                                 'issue': 'Grammar Error',
#                                 'text': text,
#                                 'corrected': match.replacements
#                             })
#     return issues

# # Decimal Consistency Validation
# def validate_decimal_consistency(slide, slide_index):
#     issues = []
#     decimal_pattern = re.compile(r'\d+[\.,]\d+')  # Pattern to match decimal numbers with either '.' or ','
#     decimal_places_set = set()

#     for shape in slide.shapes:
#         if shape.has_text_frame:
#             for paragraph in shape.text_frame.paragraphs:
#                 for run in paragraph.runs:
#                     matches = decimal_pattern.findall(run.text)
#                     for match in matches:
#                         logging.debug(f"Found decimal: {match}")  # Debugging line
#                         decimal_places = len(match.split(',')[1] if ',' in match else match.split('.')[1])
#                         decimal_places_set.add(decimal_places)

#     if len(decimal_places_set) > 1:
#         for match in matches:
#             issues.append({
#                 'slide': slide_index,
#                 'issue': 'Inconsistent Decimal Points',
#                 'text': match,
#                 'details': f'Found inconsistent decimal points: {list(decimal_places_set)}'
#             })

#     return issues
# # Highlight issues in PPT
# def highlight_ppt(input_ppt, output_ppt, issues):
#     presentation = Presentation(input_ppt)
#     for issue in issues:
#         slide_index = issue['slide'] - 1
#         slide = presentation.slides[slide_index]
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 for paragraph in shape.text_frame.paragraphs:
#                     for run in paragraph.runs:
#                         if issue['text'] in run.text:
#                             run.font.color.rgb = RGBColor(255, 255, 0)
#     presentation.save(output_ppt)

# # Save results to CSV
# def save_to_csv(issues, output_csv):
#     with open(output_csv, mode='w', newline='', encoding='utf-8') as file:
#         writer = csv.DictWriter(file, fieldnames=['slide', 'issue', 'text', 'corrected', 'details'])
#         writer.writeheader()
#         writer.writerows(issues)

# # Password Protection
# PREDEFINED_PASSWORD = "securepassword123"

# def password_protection():
#     if "authenticated" not in st.session_state:
#         st.session_state.authenticated = False
#     if not st.session_state.authenticated:
#         with st.form("password_form", clear_on_submit=True):
#             password_input = st.text_input("Enter Password", type="password")
#             submitted = st.form_submit_button("Submit")
#             if submitted and password_input == PREDEFINED_PASSWORD:
#                 st.session_state.authenticated = True
#                 st.success("Access Granted! Please click 'Submit' again to proceed.")
#             elif submitted:
#                 st.error("Incorrect Password")
#         return False
#     return True

# def main():
#     if not password_protection():
#         return

#     st.title("PPT Validator")
#     uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"])
#     font_options = ["Arial", "Calibri", "Times New Roman", "Verdana", "Helvetica", "EYInterstate"]
#     default_font = st.selectbox("Select the default font for validation", font_options)
#     validation_option = st.radio("Validation Option:", ["All Slides", "Custom Range"])

#     if uploaded_file:
#         with tempfile.TemporaryDirectory() as tmpdir:
#             temp_ppt_path = Path(tmpdir) / "uploaded_ppt.pptx"
#             with open(temp_ppt_path, "wb") as f:
#                 f.write(uploaded_file.getbuffer())

#             presentation = Presentation(temp_ppt_path)
#             total_slides = len(presentation.slides)

#             # Slide Range Selection
#             start_slide, end_slide = 1, total_slides
#             if validation_option == "Custom Range":
#                 start_slide = st.number_input("From Slide", min_value=1, max_value=total_slides, value=1)
#                 end_slide_default = min(total_slides, 100)  # Ensure default value does not exceed total slides
#                 end_slide = st.number_input("To Slide", min_value=start_slide, max_value=total_slides, value=end_slide_default)

#             if st.button("Run Validation"):
#                 progress_bar = st.progress(0)
#                 progress_text = st.empty()
#                 issues = []

#                 # Parallel Processing
#                 with ThreadPoolExecutor() as executor:
#                     futures = []
#                     for slide_index in range(start_slide - 1, end_slide):
#                         slide = presentation.slides[slide_index]
#                         futures.append(executor.submit(validate_spelling_slide, slide, slide_index + 1))
#                         futures.append(executor.submit(validate_fonts_slide, slide, slide_index + 1, default_font))
#                         futures.append(executor.submit(validate_grammar_slide, slide, slide_index + 1))
#                         futures.append(executor.submit(validate_decimal_consistency, slide, slide_index + 1))

#                     for i, future in enumerate(futures):
#                         issues.extend(future.result())
#                         progress_percent = int((i + 1) / len(futures) * 100)
#                         progress_text.text(f"Progress: {progress_percent}%")
#                         progress_bar.progress(progress_percent / 100)

#                 # Save Results
#                 csv_output_path = Path(tmpdir) / "validation_report.csv"
#                 highlighted_ppt_path = Path(tmpdir) / "highlighted_presentation.pptx"
#                 save_to_csv(issues, csv_output_path)
#                 highlight_ppt(temp_ppt_path, highlighted_ppt_path, issues)

#                 # Store results in session state
#                 st.session_state['csv_output'] = csv_output_path.read_bytes()
#                 st.session_state['ppt_output'] = highlighted_ppt_path.read_bytes()
#                 st.success("Validation completed!")

#                 # Display Download Buttons
#                 if 'csv_output' in st.session_state:
#                     st.download_button("Download Validation Report (CSV)", st.session_state['csv_output'], file_name="validation_report.csv")
#                 if 'ppt_output' in st.session_state:
#                     st.download_button("Download Highlighted PPT", st.session_state['ppt_output'], file_name="highlighted_presentation.pptx")

# if __name__ == "__main__":
#     main()


##########################################################################################

from concurrent.futures import ThreadPoolExecutor
import streamlit as st
import tempfile
from pathlib import Path
from pptx import Presentation
from spellchecker import SpellChecker
import language_tool_python
import csv
import re
import string
from pptx.dml.color import RGBColor
import logging

# Initialize LanguageTool
def initialize_language_tool():
    try:
        return language_tool_python.LanguageToolPublicAPI('en-US')
    except Exception as e:
        st.error(f"LanguageTool initialization failed: {e}")
        return None

grammar_tool = initialize_language_tool()

# Custom dictionary
TECHNICAL_TERMS = {
    "TensorFlow", "Caret", "ML", "DALL-E", "MLOps", "PyTorch", "MENA", "Keras", "Scikit-learn", "NumPy", "Pandas",
    "Matplotlib", "OpenAI", "GPT-3", "Deep Learning", "Neural Network", "Data Science", "Seaborn", "Jupyter",
    "Anaconda", "Reinforcement Learning", "Supervised Learning", "Unsupervised Learning",
    "Natural Language Processing", "Computer Vision", "Big Data", "Data Mining", "Feature Engineering",
    "Hyperparameter", "Gradient Descent", "Convolutional Neural Network", "Recurrent Neural Network",
    "Support Vector Machine", "Decision Tree", "Random Forest", "Ensemble Learning", "Clustering",
    "Dimensionality Reduction", "Principal Component Analysis", "Exploratory Data Analysis", "Model Evaluation",
    "Cross-Validation", "Overfitting", "Underfitting", "Batch Normalization", "Dropout", "Activation Function",
    "Loss Function", "Backpropagation", "Transfer Learning", "Generative Adversarial Network", "Autoencoder",
    "Tokenization", "Embedding", "Word2Vec", "BERT", "OpenCV", "Flask", "Django", "REST API", "GraphQL", "SQL",
    "NoSQL", "MongoDB", "PostgreSQL", "MySQL", "Firebase", "Cloud Computing", "AWS", "Azure", "Google Cloud",
    "Docker", "Kubernetes", "CI/CD", "DevOps", "Agile", "Scrum", "Kanban", "Git", "GitHub", "Bitbucket",
    "Version Control", "API", "SDK", "Microservices", "Blockchain", "Cryptocurrency", "IoT", "Edge Computing",
    "Quantum Computing", "Augmented Reality", "Virtual Reality", "3D Printing", "Cybersecurity",
    "Penetration Testing", "Phishing", "Malware", "Ransomware", "Firewall", "VPN", "SSL", "Encryption",
    "Decryption", "Hashing", "Digital Signature", "Data Privacy", "GDPR"
}

NUMERIC_TERMS = {f"{i}+" for i in range(1, 101)}

# Initialize SpellChecker
spell = SpellChecker()
spell.word_frequency.load_words(TECHNICAL_TERMS.union(NUMERIC_TERMS))

# Exemption function
def is_exempted(word):
    return word in TECHNICAL_TERMS or re.match(r"^\d+\+?$", word)

# Spelling Validation
def validate_spelling_slide(slide, slide_index):
    issues = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    words = re.findall(r"\b[\w+]+\b", run.text)
                    for word in words:
                        clean_word = word.strip(string.punctuation)
                        if is_exempted(clean_word):
                            continue
                        if clean_word.lower() not in spell:
                            correction = spell.correction(clean_word)
                            if correction and correction != clean_word:
                                issues.append({
                                    'slide': slide_index,
                                    'issue': 'Misspelling',
                                    'text': word,
                                    'corrected': correction
                                })
    return issues

# Font Validation
def validate_fonts_slide(slide, slide_index, default_font):
    issues = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip() and run.font.name != default_font:
                        issues.append({
                            'slide': slide_index,
                            'issue': 'Inconsistent Font',
                            'text': run.text,
                            'corrected': f"Expected: {default_font}, Found: {run.font.name}"
                        })
    return issues

# Grammar Validation
def validate_grammar_slide(slide, slide_index):
    issues = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text.strip()
                    if text and grammar_tool:
                        matches = grammar_tool.check(text)
                        for match in matches:
                            issues.append({
                                'slide': slide_index,
                                'issue': 'Grammar Error',
                                'text': text,
                                'corrected': match.replacements
                            })
    return issues

# Decimal Consistency Validation
def validate_decimal_consistency(slide, slide_index):
    issues = []
    decimal_pattern = re.compile(r'\d+[\.,]\d+')
    decimal_places_set = set()
    all_matches = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    matches = decimal_pattern.findall(run.text)
                    all_matches.extend(matches)
                    for match in matches:
                        decimal_places = len(match.split(',')[1] if ',' in match else match.split('.')[1])
                        decimal_places_set.add(decimal_places)
    if len(decimal_places_set) > 1:
        for match in all_matches:
            issues.append({
                'slide': slide_index,
                'issue': 'Inconsistent Decimal Points',
                'text': match,
                'details': f'Found inconsistent decimal points: {list(decimal_places_set)}'
            })
    return issues

# Million Notation Validation
def validate_million_notations(slide, slide_index):
    issues = []
    million_patterns = {
        r'\b\d+M\b': 'M', r'\b\d+\s?Million\b': 'Million', r'\b\d+mn\b': 'mn', r'\b\d+\sm\b': 'm',
        r'\b\d+MM\b': 'MM', r'\b\d+\s?Millions\b': 'Millions', r'\b\d+\s?Juta\b': 'Juta'
    }
    notation_set = set()
    all_matches = []
    logging.debug(f"Slide {slide_index}: Checking shapes for million notations")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for pattern, notation in million_patterns.items():
                    matches = re.findall(pattern, run.text, re.IGNORECASE)
                    all_matches.extend(matches)
                    for match in matches:
                        notation_set.add(notation)
    if len(notation_set) > 1:
        for match in all_matches:
            issues.append({
                'slide': slide_index,
                'issue': 'Inconsistent Million Notations',
                'text': match,
                'details': f'Found inconsistent million notations: [using {", ".join(notation_set)}]'
            })
    return issues
# Highlight issues in PPT
def highlight_ppt(input_ppt, output_ppt, issues):
    presentation = Presentation(input_ppt)
    for issue in issues:
        slide_index = issue['slide'] - 1
        slide = presentation.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if issue['text'] in run.text:
                            run.font.color.rgb = RGBColor(255, 255, 0)
    presentation.save(output_ppt)

# Save results to CSV
def save_to_csv(issues, output_csv):
    with open(output_csv, mode='w', newline='', encoding='utf-8') as file:
        writer = csv.DictWriter(file, fieldnames=['slide', 'issue', 'text', 'corrected', 'details'])
        writer.writeheader()
        writer.writerows(issues)

# Password Protection
PREDEFINED_PASSWORD = "securepassword123"

def password_protection():
    if "authenticated" not in st.session_state:
        st.session_state.authenticated = False
    if not st.session_state.authenticated:
        with st.form("password_form", clear_on_submit=True):
            password_input = st.text_input("Enter Password", type="password")
            submitted = st.form_submit_button("Submit")
            if submitted and password_input == PREDEFINED_PASSWORD:
                st.session_state.authenticated = True
                st.success("Access Granted! Please click 'Submit' again to proceed.")
            elif submitted:
                st.error("Incorrect Password")
        return False
    return True

def main():
    if not password_protection():
        return

    st.title("PPT Validator")
    uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"])
    font_options = ["Arial", "Calibri", "Times New Roman", "Verdana", "Helvetica", "EYInterstate"]
    default_font = st.selectbox("Select the default font for validation", font_options)
    validation_option = st.radio("Validation Option:", ["All Slides", "Custom Range"])

    if uploaded_file:
        with tempfile.TemporaryDirectory() as tmpdir:
            temp_ppt_path = Path(tmpdir) / "uploaded_ppt.pptx"
            with open(temp_ppt_path, "wb") as f:
                f.write(uploaded_file.getbuffer())

            presentation = Presentation(temp_ppt_path)
            total_slides = len(presentation.slides)

            # Slide Range Selection
            start_slide, end_slide = 1, total_slides
            if validation_option == "Custom Range":
                start_slide = st.number_input("From Slide", min_value=1, max_value=total_slides, value=1)
                end_slide_default = min(total_slides, 100)  # Ensure default value does not exceed total slides
                end_slide = st.number_input("To Slide", min_value=start_slide, max_value=total_slides, value=end_slide_default)

            if st.button("Run Validation"):
                progress_bar = st.progress(0)
                progress_text = st.empty()
                issues = []

                # Parallel Processing
                with ThreadPoolExecutor() as executor:
                    futures = []
                    for slide_index in range(start_slide - 1, end_slide):
                        slide = presentation.slides[slide_index]
                        futures.append(executor.submit(validate_spelling_slide, slide, slide_index + 1))
                        futures.append(executor.submit(validate_fonts_slide, slide, slide_index + 1, default_font))
                        futures.append(executor.submit(validate_grammar_slide, slide, slide_index + 1))
                        futures.append(executor.submit(validate_decimal_consistency, slide, slide_index + 1))
                        futures.append(executor.submit(validate_million_notations, slide, slide_index + 1))  # Added function call

                    for i, future in enumerate(futures):
                        result = future.result()
                        logging.debug(f"Validation result for future {i}: {result}")  # Added logging
                        issues.extend(result)
                        progress_percent = int((i + 1) / len(futures) * 100)
                        progress_text.text(f"Progress: {progress_percent}%")
                        progress_bar.progress(progress_percent / 100)

                # Save Results
                csv_output_path = Path(tmpdir) / "validation_report.csv"
                highlighted_ppt_path = Path(tmpdir) / "highlighted_presentation.pptx"
                save_to_csv(issues, csv_output_path)
                highlight_ppt(temp_ppt_path, highlighted_ppt_path, issues)

                # Store results in session state
                st.session_state['csv_output'] = csv_output_path.read_bytes()
                st.session_state['ppt_output'] = highlighted_ppt_path.read_bytes()
                st.success("Validation completed!")

                # Display Download Buttons
                if 'csv_output' in st.session_state:
                    st.download_button("Download Validation Report (CSV)", st.session_state['csv_output'], file_name="validation_report.csv")
                if 'ppt_output' in st.session_state:
                    st.download_button("Download Highlighted PPT", st.session_state['ppt_output'], file_name="highlighted_presentation.pptx")

if __name__ == "__main__":
    main()
