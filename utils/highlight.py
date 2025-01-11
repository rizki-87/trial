from pptx import Presentation
from pptx.dml.color import RGBColor
import csv
import logging

# Configure logging
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')

def highlight_ppt(input_ppt, output_ppt, issues):
    """
    Highlight text in the PowerPoint presentation based on the issues found.

    Parameters:
    - input_ppt: Path to the input PowerPoint file.
    - output_ppt: Path to save the highlighted PowerPoint file.
    - issues: List of issues found in the presentation.
    """
    presentation = Presentation(input_ppt)
    for issue in issues:
        if isinstance(issue, dict):
            slide_index = issue['slide'] - 1
            slide = presentation.slides[slide_index]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if issue.get('text', '') in run.text:  # Gunakan .get() untuk menghindari KeyError
                                run.font.color.rgb = RGBColor(255, 255, 0)  # Highlight in yellow

    # Save the highlighted presentation
    presentation.save(output_ppt)

def save_to_csv(issues, output_csv):
    """
    Save the validation issues to a CSV file.

    Parameters:
    - issues: List of issues found in the presentation.
    - output_csv: Path to save the CSV file.
    """
    with open(output_csv, mode='w', newline='', encoding='utf-8') as file:
        writer = csv.DictWriter(file, fieldnames=['slide', 'issue', 'text', 'corrected', 'details'])
        writer.writeheader()
        for issue in issues:
            if isinstance(issue, dict):
                # Tambahkan logging untuk memeriksa isi issue
                logging.debug(f"Issue: {issue}")  # Log isi dari issue
                writer.writerow({
                    'slide': issue.get('slide', ''),  # Gunakan .get() untuk menghindari KeyError
                    'issue': issue.get('issue', ''),
                    'text': issue.get('text', 'N/A'),  # Ganti dengan 'N/A' jika tidak ada
                    'corrected': issue.get('corrected', ''),
                    'details': issue.get('details', '')
                })
