from pptx.dml.color import RGBColor

def validate_fonts_slide(slide, slide_index, default_font):
    issues = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip() and run.font.name != default_font:
                        issues.append({
                            'slide': slide_index,
                            'issue': 'Inconsistent Font',
                            'text': run.text,
                            'corrected': f"Expected: {default_font}, Found: {run.font.name}"
                        })
    return issues
